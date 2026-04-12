#!/usr/bin/env python3
"""
能力 5: Office 文件制作测试
测试 Excel/Word/PPT 生成
"""
import os
import sys

def test_excel_generation():
    """测试 5.1: Excel 生成"""
    print("📊 测试 Excel 生成...")
    try:
        import xlsxwriter
        
        test_file = "/tmp/output-test.xlsx"
        workbook = xlsxwriter.Workbook(test_file)
        worksheet = workbook.add_worksheet("Monthly Report")
        
        # Add formatting
        bold = workbook.add_format({'bold': True, 'bg_color': '#4472C4', 'font_color': 'white'})
        currency = workbook.add_format({'num_format': '#,##0.00'})
        
        # Headers
        headers = ["Month", "Revenue", "Cost", "Profit"]
        for col, header in enumerate(headers):
            worksheet.write(0, col, header, bold)
        
        # Data
        data = [
            ["Jan", 100000, 60000],
            ["Feb", 120000, 70000],
            ["Mar", 150000, 80000],
        ]
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, value in enumerate(row_data):
                worksheet.write(row_idx, col_idx, value)
            # Formula for profit
            worksheet.write_formula(row_idx, 3, f'=B{row_idx+1}-C{row_idx+1}', currency)
        
        # Total row
        worksheet.write(4, 0, "Total", bold)
        worksheet.write_formula(4, 1, '=SUM(B2:B4)', currency)
        worksheet.write_formula(4, 2, '=SUM(C2:C4)', currency)
        worksheet.write_formula(4, 3, '=B5-C5', currency)
        
        worksheet.set_column(0, 0, 10)
        worksheet.set_column(1, 3, 15)
        workbook.close()
        
        # Verify
        assert os.path.exists(test_file)
        size = os.path.getsize(test_file)
        print(f"  生成文件: {test_file}")
        print(f"  文件大小: {size:,} bytes")
        print("  ✅ Excel 生成通过")
        return True
    except Exception as e:
        print(f"  ❌ Excel 生成失败: {e}")
        return False

def test_word_generation():
    """测试 5.2: Word 生成"""
    print("\n📝 测试 Word 生成...")
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # Title
        title = doc.add_heading('季度经营报告', level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Sections
        doc.add_heading('一、概述', level=1)
        doc.add_paragraph('本季度整体经营状况良好，收入同比增长 15%。')
        
        doc.add_heading('二、财务数据', level=1)
        # Add table
        table = doc.add_table(rows=4, cols=3, style='Light Shading Accent 1')
        headers = ['指标', 'Q1', 'Q2']
        for i, header in enumerate(headers):
            table.rows[0].cells[i].text = header
        data = [
            ['收入 (万)', '500', '650'],
            ['成本 (万)', '300', '380'],
            ['利润 (万)', '200', '270'],
        ]
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, value in enumerate(row_data):
                table.rows[row_idx].cells[col_idx].text = value
        
        doc.add_heading('三、总结', level=1)
        doc.add_paragraph('预计下季度继续保持增长态势。')
        
        test_file = "/tmp/output-test.docx"
        doc.save(test_file)
        
        # Verify
        assert os.path.exists(test_file)
        size = os.path.getsize(test_file)
        print(f"  生成文件: {test_file}")
        print(f"  文件大小: {size:,} bytes")
        print("  ✅ Word 生成通过")
        return True
    except Exception as e:
        print(f"  ❌ Word 生成失败: {e}")
        return False

def test_ppt_generation():
    """测试 5.3: PPT 生成"""
    print("\n📊 测试 PPT 生成...")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        
        # Title slide
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        slide.shapes.title.text = "季度汇报"
        slide.placeholders[1].text = "2026年Q2"
        
        # Content slide
        content_slide = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide)
        slide.shapes.title.text = "经营概况"
        
        # Add bullet points
        tf = slide.placeholders[1].text_frame
        tf.text = "收入同比增长 15%"
        p = tf.add_paragraph()
        p.text = "成本控制在预算内"
        p = tf.add_paragraph()
        p.text = "利润率提升 3%"
        
        # Table slide
        blank_slide = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide)
        slide.shapes.title.text = "财务数据"
        
        table = slide.shapes.add_table(4, 3, Inches(1), Inches(2), Inches(8), Inches(3)).table
        table.cell(0, 0).text = "指标"
        table.cell(0, 1).text = "Q1"
        table.cell(0, 2).text = "Q2"
        data = [["收入", "500万", "650万"], ["成本", "300万", "380万"], ["利润", "200万", "270万"]]
        for i, row_data in enumerate(data, start=1):
            for j, val in enumerate(row_data):
                table.cell(i, j).text = val
        
        test_file = "/tmp/output-test.pptx"
        prs.save(test_file)
        
        # Verify
        assert os.path.exists(test_file)
        size = os.path.getsize(test_file)
        print(f"  生成文件: {test_file}")
        print(f"  文件大小: {size:,} bytes")
        print("  ✅ PPT 生成通过")
        return True
    except Exception as e:
        print(f"  ❌ PPT 生成失败: {e}")
        return False

def main():
    print("=" * 50)
    print("能力 5: Office 文件制作测试")
    print("=" * 50)
    
    results = {}
    results["test_5_1_excel"] = "✅ 通过" if test_excel_generation() else "❌ 失败"
    results["test_5_2_word"] = "✅ 通过" if test_word_generation() else "❌ 失败"
    results["test_5_3_ppt"] = "✅ 通过" if test_ppt_generation() else "❌ 失败"
    
    print("\n" + "=" * 50)
    print("测试结果汇总:")
    for k, v in results.items():
        print(f"  {k}: {v}")
    passed = sum(1 for v in results.values() if "✅" in v)
    print(f"  通过率: {passed}/{len(results)}")
    print("=" * 50)
    
    return passed == len(results)

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
